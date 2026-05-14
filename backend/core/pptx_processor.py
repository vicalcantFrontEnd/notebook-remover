"""
PPTX watermark removal processor.
Instead of inpainting, this removes shapes/elements whose bounding box
overlaps with the selected region. Preserves the editable PPTX format.
"""

import zipfile
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
import cv2
import numpy as np


def get_pptx_info(path: str) -> dict:
    """Return slide width, height (in pixels at 96 DPI), and slide count."""
    prs = Presentation(str(Path(path).resolve()))
    # Convert EMUs to pixels at 96 DPI (1 inch = 914400 EMU, 96 px)
    emu_per_px = 914400 / 96
    width = int(prs.slide_width / emu_per_px)
    height = int(prs.slide_height / emu_per_px)
    slide_count = len(prs.slides)
    return {
        "width": width,
        "height": height,
        "slide_count": slide_count,
    }


def extract_first_slide_image(path: str) -> np.ndarray:
    """
    Extract thumbnail from PPTX zip (docProps/thumbnail.jpeg).
    If not available, create a white canvas with slide dimensions.
    """
    try:
        with zipfile.ZipFile(str(Path(path).resolve()), "r") as z:
            # Try common thumbnail paths
            for thumb_path in ["docProps/thumbnail.jpeg", "docProps/thumbnail.png"]:
                if thumb_path in z.namelist():
                    thumb_data = z.read(thumb_path)
                    arr = np.frombuffer(thumb_data, dtype=np.uint8)
                    img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
                    if img is not None:
                        return img
    except (zipfile.BadZipFile, KeyError):
        pass

    # Fallback: white canvas
    info = get_pptx_info(path)
    canvas = np.ones((info["height"], info["width"], 3), dtype=np.uint8) * 255
    return canvas


def _regions_overlap(
    shape_box: tuple[int, int, int, int],
    region: tuple[int, int, int, int],
    threshold: float = 0.3,
) -> bool:
    """
    Check if a shape's bounding box overlaps significantly with the region.
    Both are (x, y, w, h) in the same coordinate space.
    Returns True if intersection area / shape area >= threshold.
    """
    sx, sy, sw, sh = shape_box
    rx, ry, rw, rh = region

    # Intersection
    ix1 = max(sx, rx)
    iy1 = max(sy, ry)
    ix2 = min(sx + sw, rx + rw)
    iy2 = min(sy + sh, ry + rh)

    if ix2 <= ix1 or iy2 <= iy1:
        return False

    intersection = (ix2 - ix1) * (iy2 - iy1)
    shape_area = sw * sh
    if shape_area == 0:
        return False

    return (intersection / shape_area) >= threshold


def process_pptx(
    input_path: str,
    output_path: str,
    region: tuple[int, int, int, int],
    progress_callback=None,
    cancel_flag=None,
) -> None:
    """
    Remove shapes from all slides that overlap with the selected region.

    Region coordinates are in pixels (at 96 DPI). They are converted to
    EMUs for comparison with shape positions.

    Args:
        region: (x, y, w, h) in pixel coordinates.
    """
    prs = Presentation(str(input_path))

    # Convert pixel region to EMUs
    emu_per_px = 914400 / 96
    region_emu = (
        int(region[0] * emu_per_px),
        int(region[1] * emu_per_px),
        int(region[2] * emu_per_px),
        int(region[3] * emu_per_px),
    )

    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        if cancel_flag and cancel_flag():
            return

        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            shape_box = (
                int(shape.left),
                int(shape.top),
                int(shape.width),
                int(shape.height),
            )
            if _regions_overlap(shape_box, region_emu):
                shapes_to_remove.append(shape)

        # Remove shapes from slide's XML
        for shape in shapes_to_remove:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

        if progress_callback:
            progress_callback(i + 1, total_slides)

    prs.save(str(Path(output_path).resolve()))
